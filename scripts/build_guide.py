from __future__ import annotations

import re
import shutil
import subprocess
from dataclasses import dataclass
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.util import Inches, Pt
from reportlab.lib import colors
from reportlab.lib.enums import TA_CENTER, TA_LEFT
from reportlab.lib.pagesizes import landscape, A4
from reportlab.lib.styles import ParagraphStyle
from reportlab.lib.units import mm
from reportlab.graphics.shapes import Circle, Drawing, Line, Rect, String
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
from reportlab.platypus import Image as PdfImage
from reportlab.platypus import PageBreak, Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle

from guidebook_common import (
    guide_spots,
    meal_recommendations,
    route_map_points,
    route_map_url,
    route_points,
    short_text,
    today_theme,
    todays_tips,
)


ROOT = Path(__file__).resolve().parents[1]
ITINERARY_DIR = ROOT / "itinerary" / "days"
IMAGES_DIR = ROOT / "images"
OUTPUT_DIR = ROOT / "output"
PPTX_PATH = OUTPUT_DIR / "hokkaido-family-travel-guide.pptx"
PDF_PATH = OUTPUT_DIR / "hokkaido-family-travel-guide.pdf"
PLACEHOLDER_DIR = OUTPUT_DIR / "_placeholders"
COVER_IMAGE_PATH = IMAGES_DIR / "cover.png"

JP_FONT = "Yu Gothic"
EN_FONT = "Aptos"

NAVY = RGBColor(20, 55, 83)
BLUE = RGBColor(44, 115, 160)
SKY = RGBColor(226, 241, 248)
SAND = RGBColor(245, 238, 220)
CORAL = RGBColor(216, 99, 82)
INK = RGBColor(34, 43, 52)
MUTED = RGBColor(102, 112, 122)
WHITE = RGBColor(255, 255, 255)
PALE = RGBColor(248, 250, 252)


@dataclass
class TimelineItem:
    time: str
    type: str
    place: str
    detail: str
    duration: str


@dataclass
class Restaurant:
    meal: str
    name: str
    area: str
    memo: str


@dataclass
class PhotoSpot:
    place: str
    image: str
    caption: str
    credit: str


@dataclass
class DayPlan:
    date: str
    day: str
    title: str
    area: str
    hero: str
    summary: str
    timeline: list[TimelineItem]
    photos: list[PhotoSpot]
    restaurants: list[Restaurant]
    notes: list[str]


def clean_cell(value: str) -> str:
    return value.strip().replace("<br>", "\n")


def parse_table(lines: list[str]) -> list[list[str]]:
    rows: list[list[str]] = []
    for line in lines:
        line = line.strip()
        if not line.startswith("|"):
            continue
        cells = [clean_cell(c) for c in line.strip("|").split("|")]
        if all(set(c) <= {"-", ":", " "} for c in cells):
            continue
        rows.append(cells)
    return rows


def section(text: str, name: str) -> str:
    pattern = rf"^## {re.escape(name)}\s*$"
    matches = list(re.finditer(pattern, text, flags=re.MULTILINE))
    if not matches:
        return ""
    start = matches[0].end()
    next_match = re.search(r"^## .+$", text[start:], flags=re.MULTILINE)
    end = start + next_match.start() if next_match else len(text)
    return text[start:end].strip()


def parse_day(path: Path) -> DayPlan:
    raw = path.read_text(encoding="utf-8")
    meta: dict[str, str] = {}
    for line in raw.splitlines()[1:20]:
        if not line.strip():
            continue
        if line.startswith("## "):
            break
        if ":" in line:
            key, value = line.split(":", 1)
            meta[key.strip()] = value.strip()

    summary = section(raw, "Summary").replace("\n", " ").strip()

    timeline_rows = parse_table(section(raw, "Timeline").splitlines())
    timeline = [
        TimelineItem(*row[:5])
        for row in timeline_rows[1:]
        if len(row) >= 5
    ]

    restaurant_rows = parse_table(section(raw, "Restaurants").splitlines())
    restaurants = [
        Restaurant(*row[:4])
        for row in restaurant_rows[1:]
        if len(row) >= 4
    ]

    photo_rows = parse_table(section(raw, "Photos").splitlines())
    photos = [
        PhotoSpot(*row[:4])
        for row in photo_rows[1:]
        if len(row) >= 4
    ]

    notes = [
        line.strip()[2:].strip()
        for line in section(raw, "Notes").splitlines()
        if line.strip().startswith("- ")
    ]

    return DayPlan(
        date=meta.get("date", path.stem),
        day=meta.get("day", ""),
        title=meta.get("title", path.stem),
        area=meta.get("area", ""),
        hero=meta.get("hero", meta.get("title", path.stem)),
        summary=summary,
        timeline=timeline,
        photos=photos,
        restaurants=restaurants,
        notes=notes,
    )


def load_days() -> list[DayPlan]:
    return [parse_day(path) for path in sorted(ITINERARY_DIR.glob("*.md"))]


def add_textbox(slide, x, y, w, h, text, font_size=18, color=INK, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = JP_FONT
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_label(slide, x, y, text, fill=CORAL, color=WHITE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(1.45), Inches(0.36))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    tf = shape.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = JP_FONT
    r.font.size = Pt(11)
    r.font.bold = True
    r.font.color.rgb = color
    return shape


def image_candidates(day: DayPlan) -> list[Path]:
    names = [day.date, day.hero, day.area, day.title]
    normalized: list[str] = []
    for name in names:
        for part in re.split(r"[/／ ]+", name):
            if part:
                normalized.append(part)
    candidates: list[Path] = []
    for name in names + normalized:
        for ext in (".jpg", ".jpeg", ".png", ".webp"):
            candidates.append(IMAGES_DIR / f"{name}{ext}")
    return candidates


def make_placeholder(label: str, path: Path, size=(1600, 900)) -> Path:
    path.parent.mkdir(parents=True, exist_ok=True)
    img = Image.new("RGB", size, (228, 241, 248))
    draw = ImageDraw.Draw(img)
    for i in range(0, size[0], 80):
        draw.line((i, 0, i - 600, size[1]), fill=(210, 230, 240), width=6)
    try:
        font_big = ImageFont.truetype("/System/Library/Fonts/Supplemental/Arial Unicode.ttf", 84)
        font_small = ImageFont.truetype("/System/Library/Fonts/Supplemental/Arial Unicode.ttf", 34)
    except Exception:
        font_big = ImageFont.load_default()
        font_small = ImageFont.load_default()
    title = label
    subtitle = "写真準備中"
    tw = draw.textbbox((0, 0), title, font=font_big)
    sw = draw.textbbox((0, 0), subtitle, font=font_small)
    draw.rounded_rectangle((120, 260, 1480, 640), radius=36, fill=(255, 255, 255), outline=(44, 115, 160), width=4)
    draw.text(((size[0] - (tw[2] - tw[0])) / 2, 350), title, fill=(20, 55, 83), font=font_big)
    draw.text(((size[0] - (sw[2] - sw[0])) / 2, 470), subtitle, fill=(102, 112, 122), font=font_small)
    img.save(path)
    return path


def hero_image(day: DayPlan) -> Path:
    for path in image_candidates(day):
        if path.exists():
            return path
    safe = re.sub(r"[^0-9A-Za-zぁ-んァ-ヶ一-龠ー_-]+", "_", day.hero)[:50] or day.date
    return make_placeholder(day.hero, PLACEHOLDER_DIR / f"{day.date}_{safe}.png")


def cover_image(days: list[DayPlan]) -> Path:
    if COVER_IMAGE_PATH.exists():
        return COVER_IMAGE_PATH
    return hero_image(days[0])


def photo_image(day: DayPlan, photo: PhotoSpot) -> Path:
    if photo.image:
        explicit = IMAGES_DIR / photo.image
        if explicit.exists():
            return explicit
    safe = re.sub(r"[^0-9A-Za-zぁ-んァ-ヶ一-龠ー_-]+", "_", photo.place)[:50] or day.date
    return make_placeholder(photo.place, PLACEHOLDER_DIR / f"{day.date}_{safe}.png")


def spot_image(day: DayPlan, spot: dict[str, str]) -> Path:
    image_name = spot.get("image", "")
    if image_name:
        explicit = IMAGES_DIR / image_name
        if explicit.exists():
            return explicit
    label = spot.get("place") or day.hero
    safe = re.sub(r"[^0-9A-Za-zぁ-んァ-ヶ一-龠ー_-]+", "_", label)[:50] or day.date
    return make_placeholder(label, PLACEHOLDER_DIR / f"{day.date}_{safe}.png")


def add_picture_cover(slide, image_path: Path, x, y, w, h):
    slide.shapes.add_picture(str(image_path), x, y, width=w, height=h)


def style_background(slide, color=RGBColor(255, 255, 255)):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_timeline(slide, items: list[TimelineItem], start_y=Inches(1.55), max_items=8, compact=False):
    y = start_y
    card_h = Inches(0.49 if compact else 0.55)
    step = Inches(0.54 if compact else 0.62)
    time_size = 8.8 if compact else 9.5
    place_size = 8.4 if compact else 9.0
    detail_size = 7.6 if compact else 8.3
    for item in items[:max_items]:
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.55), y, Inches(7.0), card_h)
        card.fill.solid()
        card.fill.fore_color.rgb = PALE
        card.line.color.rgb = RGBColor(222, 228, 234)
        add_textbox(slide, Inches(0.68), y + Inches(0.07), Inches(0.9), Inches(0.25), item.time, time_size, BLUE, True, PP_ALIGN.CENTER)
        add_textbox(slide, Inches(1.55), y + Inches(0.07), Inches(0.75), Inches(0.25), item.type, time_size, CORAL, True, PP_ALIGN.CENTER)
        add_textbox(slide, Inches(2.28), y + Inches(0.06), Inches(1.35), Inches(0.28), item.place, place_size, NAVY, True)
        add_textbox(slide, Inches(3.55), y + Inches(0.05), Inches(3.15), Inches(0.32), short_text(item.detail, 72), detail_size, INK)
        add_textbox(slide, Inches(6.72), y + Inches(0.07), Inches(0.70), Inches(0.25), item.duration, detail_size, MUTED, False, PP_ALIGN.RIGHT)
        y += step


def add_card(slide, x, y, w, h, fill=WHITE, line=RGBColor(218, 226, 232)):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    card.fill.solid()
    card.fill.fore_color.rgb = fill
    card.line.color.rgb = line
    return card


def add_guide_heading(slide, x, y, text, size=13):
    add_textbox(slide, x, y, Inches(2.2), Inches(0.26), text, size, BLUE, True)


def add_route_map(slide, day: DayPlan, x, y, w, h):
    add_card(slide, x, y, w, h, RGBColor(255, 252, 246), RGBColor(232, 217, 196))
    add_guide_heading(slide, x + Inches(0.18), y + Inches(0.14), "Today's Map", 13)
    link = add_textbox(slide, x + w - Inches(1.28), y + Inches(0.14), Inches(1.05), Inches(0.22), "Google Map", 7.2, BLUE, True, PP_ALIGN.RIGHT)
    link.text_frame.paragraphs[0].runs[0].hyperlink.address = route_map_url(day)

    points = route_map_points(day, max_points=8)
    if not points:
        return

    map_x = x + Inches(0.18)
    map_y = y + Inches(0.52)
    map_w = int(w * 0.47)
    map_h = h - Inches(0.72)
    sketch = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, map_x, map_y, map_w, map_h)
    sketch.fill.solid()
    sketch.fill.fore_color.rgb = SKY
    sketch.line.color.rgb = RGBColor(203, 223, 232)

    centers: list[tuple[int, int]] = []
    for index, point in enumerate(points):
        cx = int(map_x + map_w * (float(point["x"]) / 100))
        cy = int(map_y + map_h * (float(point["y"]) / 100))
        centers.append((cx, cy))
        if index > 0:
            x1, y1 = centers[index - 1]
            connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, cx, cy)
            connector.line.color.rgb = RGBColor(44, 115, 160)
            connector.line.width = Pt(2.0)

    for index, point in enumerate(points):
        cx, cy = centers[index]
        fill = CORAL if index in {0, len(points) - 1} else WHITE
        text_color = WHITE if index in {0, len(points) - 1} else BLUE
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx - Inches(0.11), cy - Inches(0.11), Inches(0.22), Inches(0.22))
        dot.fill.solid()
        dot.fill.fore_color.rgb = fill
        dot.line.color.rgb = BLUE if fill == WHITE else CORAL
        dot.text_frame.clear()
        dot.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = dot.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = str(index + 1)
        r.font.name = EN_FONT
        r.font.size = Pt(6.5)
        r.font.bold = True
        r.font.color.rgb = text_color

    list_x = map_x + map_w + Inches(0.20)
    list_w = x + w - list_x - Inches(0.18)
    row_h = min(Inches(0.36), (map_h - Inches(0.12)) / max(len(points), 1))
    for index, point in enumerate(points):
        py = map_y + Inches(0.05) + row_h * index
        add_textbox(slide, list_x, py, Inches(0.22), Inches(0.18), f"{index + 1}", 6.4, CORAL, True, PP_ALIGN.CENTER)
        add_textbox(slide, list_x + Inches(0.24), py, list_w - Inches(0.26), Inches(0.18), short_text(point["place"], 18), 7.2, NAVY, True)
        if point.get("leg") and index < len(points) - 1:
            add_textbox(slide, list_x + Inches(0.28), py + Inches(0.17), list_w - Inches(0.30), Inches(0.14), point["leg"], 5.6, MUTED)


def add_theme_card(slide, day: DayPlan, x, y, w, h):
    add_card(slide, x, y, w, h, RGBColor(247, 251, 250), RGBColor(210, 229, 226))
    add_guide_heading(slide, x + Inches(0.18), y + Inches(0.14), "Today's Theme", 13)
    add_textbox(slide, x + Inches(0.18), y + Inches(0.52), w - Inches(0.36), h - Inches(0.65), today_theme(day), 10.0, INK)


def add_meal_recommendations(slide, day: DayPlan, x, y, w, h):
    add_card(slide, x, y, w, h, RGBColor(255, 250, 240), RGBColor(234, 221, 199))
    add_guide_heading(slide, x + Inches(0.16), y + Inches(0.12), "Food Picks", 12)
    recs = meal_recommendations(day)
    item_h = (h - Inches(0.45)) / max(len(recs), 1)
    for index, rec in enumerate(recs[:2]):
        iy = y + Inches(0.48) + item_h * index
        add_textbox(slide, x + Inches(0.18), iy, Inches(1.3), Inches(0.18), rec["label"], 7.6, CORAL, True)
        add_textbox(slide, x + Inches(1.25), iy, w - Inches(1.42), Inches(0.22), rec["stars"], 7.6, RGBColor(198, 122, 36), True, PP_ALIGN.RIGHT)
        add_textbox(slide, x + Inches(0.18), iy + Inches(0.22), w - Inches(0.36), Inches(0.22), rec["name"], 8.7, NAVY, True)
        add_textbox(slide, x + Inches(0.18), iy + Inches(0.47), w - Inches(0.36), Inches(0.26), f"{rec['budget']} / {rec['popular']}", 6.8, MUTED)


def add_tips_card(slide, day: DayPlan, x, y, w, h):
    add_card(slide, x, y, w, h, RGBColor(245, 250, 240), RGBColor(218, 231, 204))
    add_guide_heading(slide, x + Inches(0.16), y + Inches(0.12), "Today's Tips", 12)
    tips = todays_tips(day, max_items=2)
    ty = y + Inches(0.48)
    for tip in tips:
        add_textbox(slide, x + Inches(0.18), ty, w - Inches(0.36), Inches(0.25), f"・{short_text(tip, 42)}", 7.0, INK)
        ty += Inches(0.34)


def add_restaurants(slide, restaurants: list[Restaurant], x=Inches(7.82), y=Inches(4.65), max_items=4):
    add_textbox(slide, x, y, Inches(4.75), Inches(0.32), "レストラン候補", 15, NAVY, True)
    y += Inches(0.42)
    for r in restaurants[:max_items]:
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(4.75), Inches(0.48))
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(255, 255, 255)
        card.line.color.rgb = RGBColor(226, 232, 240)
        add_textbox(slide, x + Inches(0.08), y + Inches(0.06), Inches(0.62), Inches(0.25), r.meal, 8.5, CORAL, True)
        add_textbox(slide, x + Inches(0.70), y + Inches(0.05), Inches(1.70), Inches(0.25), r.name, 8.5, NAVY, True)
        add_textbox(slide, x + Inches(2.35), y + Inches(0.05), Inches(2.15), Inches(0.25), r.memo, 7.5, MUTED)
        y += Inches(0.55)


def add_photo_spots(slide, day: DayPlan, x=Inches(7.82), y=Inches(0.35), max_items=3):
    if not day.photos:
        return
    add_textbox(slide, x, y, Inches(4.75), Inches(0.32), "写真で見る立ち寄り名所", 15, NAVY, True)
    y += Inches(0.42)
    photo_w = Inches(1.46)
    gap = Inches(0.18)
    for index, photo in enumerate(day.photos[:max_items]):
        px = x + index * (photo_w + gap)
        add_picture_cover(slide, photo_image(day, photo), px, y, photo_w, Inches(1.02))
        add_textbox(slide, px, y + Inches(1.05), photo_w, Inches(0.24), photo.place, 7.8, NAVY, True, PP_ALIGN.CENTER)
        add_textbox(slide, px, y + Inches(1.28), photo_w, Inches(0.34), photo.caption, 6.2, MUTED, False, PP_ALIGN.CENTER)


def add_spot_strip(slide, day: DayPlan, x, y, w, h, max_items=2):
    spots = guide_spots(day, max_items=max_items)
    if not spots:
        return
    add_guide_heading(slide, x, y, "より道スポット", 12)
    y += Inches(0.32)
    gap = Inches(0.16)
    item_w = (w - gap * (len(spots) - 1)) / len(spots)
    for index, spot in enumerate(spots):
        sx = x + (item_w + gap) * index
        add_picture_cover(slide, spot_image(day, spot), sx, y, item_w, h - Inches(0.58))
        add_textbox(slide, sx, y + h - Inches(0.50), item_w, Inches(0.18), short_text(spot["place"], 15), 6.7, NAVY, True, PP_ALIGN.CENTER)
        add_textbox(slide, sx, y + h - Inches(0.30), item_w, Inches(0.20), f"滞在 {spot['stay']} / Mapあり", 5.7, MUTED, False, PP_ALIGN.CENTER)


def chunks(items: list[TimelineItem], size: int) -> list[list[TimelineItem]]:
    return [items[index:index + size] for index in range(0, len(items), size)]


def build_pptx(days: list[DayPlan]):
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    cover = prs.slides.add_slide(blank)
    style_background(cover, SAND)
    cover_img = cover_image(days)
    add_picture_cover(cover, cover_img, Inches(6.28), Inches(0.68), Inches(6.35), Inches(4.23))
    add_textbox(cover, Inches(0.7), Inches(0.62), Inches(3.0), Inches(0.35), "HOKKAIDO FAMILY TRAVEL GUIDE", 12, BLUE, True)
    add_textbox(cover, Inches(0.68), Inches(1.45), Inches(6.1), Inches(1.0), "北海道家族旅行\nガイドブック", 33, NAVY, True)
    add_textbox(cover, Inches(0.78), Inches(3.1), Inches(5.2), Inches(0.48), "2026年7月31日 - 8月12日", 18, INK, True)
    add_textbox(cover, Inches(0.78), Inches(3.72), Inches(5.6), Inches(0.75), "仙台からフェリーで北海道へ。洞爺湖、札幌、層雲峡、道東方面をめぐる家族旅行。", 15, INK)
    add_label(cover, Inches(0.78), Inches(4.78), "2026夏版")

    for day in days:
        slide = prs.slides.add_slide(blank)
        style_background(slide, RGBColor(255, 252, 246))
        add_picture_cover(slide, hero_image(day), Inches(7.05), Inches(0.42), Inches(5.75), Inches(3.05))
        add_label(slide, Inches(0.55), Inches(0.36), f"DAY {day.day}")
        add_textbox(slide, Inches(0.55), Inches(0.78), Inches(6.15), Inches(0.58), day.title, 22, NAVY, True)
        add_textbox(slide, Inches(0.58), Inches(1.34), Inches(6.2), Inches(0.30), f"{day.date}  {day.area}", 11.0, BLUE, True)
        add_theme_card(slide, day, Inches(0.55), Inches(1.78), Inches(6.1), Inches(1.02))
        add_route_map(slide, day, Inches(0.55), Inches(3.02), Inches(6.15), Inches(4.02))
        add_meal_recommendations(slide, day, Inches(7.05), Inches(3.72), Inches(2.75), Inches(1.48))
        add_tips_card(slide, day, Inches(10.05), Inches(3.72), Inches(2.75), Inches(1.48))
        add_spot_strip(slide, day, Inches(7.05), Inches(5.58), Inches(5.75), Inches(1.35), max_items=2)

        for page_index, timeline_items in enumerate(chunks(day.timeline, 10), start=1):
            slide2 = prs.slides.add_slide(blank)
            style_background(slide2, RGBColor(252, 252, 250))
            add_label(slide2, Inches(0.55), Inches(0.35), f"DAY {day.day}")
            title = "時刻表" if page_index == 1 else "時刻表 続き"
            add_textbox(slide2, Inches(0.55), Inches(0.78), Inches(6.8), Inches(0.42), title, 22, NAVY, True)
            add_timeline(slide2, timeline_items, start_y=Inches(1.45), max_items=10, compact=True)
            add_picture_cover(slide2, hero_image(day), Inches(8.05), Inches(0.55), Inches(4.55), Inches(2.35))
            add_route_map(slide2, day, Inches(8.05), Inches(3.14), Inches(4.55), Inches(3.85))

    prs.save(PPTX_PATH)


def register_pdf_font() -> str:
    candidates = [
        "/System/Library/Fonts/Supplemental/Arial Unicode.ttf",
        "/System/Library/Fonts/ヒラギノ角ゴシック W3.ttc",
        "/System/Library/Fonts/Hiragino Sans GB.ttc",
    ]
    for candidate in candidates:
        if Path(candidate).exists():
            pdfmetrics.registerFont(TTFont("GuideJP", candidate))
            return "GuideJP"
    return "Helvetica"


def build_pdf(days: list[DayPlan]):
    font = register_pdf_font()
    page_size = landscape(A4)
    doc = SimpleDocTemplate(
        str(PDF_PATH),
        pagesize=page_size,
        leftMargin=14 * mm,
        rightMargin=14 * mm,
        topMargin=12 * mm,
        bottomMargin=10 * mm,
        title="北海道家族旅行ガイド",
    )
    base = ParagraphStyle("base", fontName=font, fontSize=9.2, leading=12.2, textColor=colors.HexColor("#222B34"), wordWrap="CJK")
    h1 = ParagraphStyle("h1", parent=base, fontSize=20, leading=25, textColor=colors.HexColor("#143753"), spaceAfter=6)
    h2 = ParagraphStyle("h2", parent=base, fontSize=13, leading=17, textColor=colors.HexColor("#2C73A0"), spaceBefore=8, spaceAfter=5)
    center = ParagraphStyle("center", parent=base, alignment=TA_CENTER)

    def p(text, style=base):
        return Paragraph(str(text).replace("&", "&amp;"), style)

    def photo_caption(photo: PhotoSpot):
        return p(f"{photo.place}<br/>{photo.caption}<br/>{photo.credit}", center)

    def pdf_table(rows, widths, header=colors.HexColor("#E2F1F8"), size=7.8):
        body = ParagraphStyle(f"tbl{size}", parent=base, fontSize=size, leading=size + 2.6, wordWrap="CJK")
        data = [[Paragraph(str(c).replace("&", "&amp;"), body) for c in row] for row in rows]
        tbl = Table(data, colWidths=widths, repeatRows=1)
        tbl.setStyle(TableStyle([
            ("FONTNAME", (0, 0), (-1, -1), font),
            ("BACKGROUND", (0, 0), (-1, 0), header),
            ("GRID", (0, 0), (-1, -1), 0.4, colors.HexColor("#B8C2CC")),
            ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
            ("LEFTPADDING", (0, 0), (-1, -1), 4),
            ("RIGHTPADDING", (0, 0), (-1, -1), 4),
            ("TOPPADDING", (0, 0), (-1, -1), 4),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 4),
        ]))
        return tbl

    def route_drawing(day: DayPlan):
        points = route_map_points(day, max_points=8)
        width = 92 * mm
        height = 58 * mm
        drawing = Drawing(width, height)
        drawing.add(Rect(0, 0, width, height, rx=8, ry=8, fillColor=colors.HexColor("#E7F1F4"), strokeColor=colors.HexColor("#CBDDE8")))
        centers = []
        for index, point in enumerate(points):
            cx = width * (float(point["x"]) / 100)
            cy = height * (1 - (float(point["y"]) / 100))
            centers.append((cx, cy))
            if index > 0:
                px, py = centers[index - 1]
                drawing.add(Line(px, py, cx, cy, strokeColor=colors.HexColor("#2C73A0"), strokeWidth=1.5))
        for index, point in enumerate(points):
            cx, cy = centers[index]
            is_edge = index in {0, len(points) - 1}
            fill = colors.HexColor("#D86352") if is_edge else colors.white
            text = colors.white if is_edge else colors.HexColor("#2C73A0")
            drawing.add(Circle(cx, cy, 5.5, fillColor=fill, strokeColor=colors.HexColor("#2C73A0"), strokeWidth=1.2))
            drawing.add(String(cx - 2.4, cy - 2.4, str(index + 1), fontName="Helvetica-Bold", fontSize=5.2, fillColor=text))
        return drawing

    def route_block(day: DayPlan):
        points = route_points(day, max_points=8)
        route_style = ParagraphStyle("route", parent=base, fontSize=7.6, leading=10.2, wordWrap="CJK")
        rows = []
        for index, point in enumerate(points):
            leg = f" / {point.get('leg', '')}" if point.get("leg") and index < len(points) - 1 else ""
            rows.append([Paragraph(f"<b>{index + 1}. {short_text(point['place'], 24)}</b>{leg}", route_style)])
        rows.append([Paragraph(f'<link href="{route_map_url(day)}">Google Mapでルートを開く</link>', route_style)])
        tbl = Table(rows, colWidths=[136 * mm])
        tbl.setStyle(TableStyle([
            ("LEFTPADDING", (0, 0), (-1, -1), 4),
            ("RIGHTPADDING", (0, 0), (-1, -1), 4),
            ("TOPPADDING", (0, 0), (-1, -1), 2),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 2),
        ]))
        return tbl

    story = [
        Spacer(1, 28),
        Paragraph("北海道家族旅行ガイドブック", ParagraphStyle("title", parent=h1, alignment=TA_CENTER, fontSize=30, leading=36)),
        Paragraph("2026年7月31日 - 8月12日", ParagraphStyle("sub", parent=center, fontSize=14, leading=18, textColor=colors.HexColor("#66707A"))),
        Spacer(1, 18),
        PdfImage(str(cover_image(days)), width=210 * mm, height=140 * mm),
        Spacer(1, 10),
        p("仙台からフェリーで北海道へ。洞爺湖、札幌、層雲峡、道東方面をめぐる家族旅行。", center),
        PageBreak(),
    ]

    for day in days:
        rec_rows = [
            [rec["label"], rec["name"], rec["stars"], rec["budget"], rec["popular"]]
            for rec in meal_recommendations(day)[:2]
        ]
        tip_rows = [[tip] for tip in todays_tips(day)]
        spot_rows = [
            [spot["place"], spot["stay"], spot["parking"], f'<link href="{spot["map_url"]}">Google Map</link>']
            for spot in guide_spots(day, max_items=2)
        ]
        header = Table(
            [
                [
                    [
                        Paragraph(f"DAY {day.day}  {day.date}", h2),
                        Paragraph(day.title, h1),
                        p(day.area),
                    ],
                    PdfImage(str(hero_image(day)), width=92 * mm, height=52 * mm),
                ]
            ],
            colWidths=[146 * mm, 96 * mm],
        )
        header.setStyle(TableStyle([
            ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
            ("LEFTPADDING", (0, 0), (-1, -1), 0),
            ("RIGHTPADDING", (0, 0), (-1, -1), 0),
        ]))
        story += [
            header,
            Spacer(1, 4),
            Paragraph("Today's Map", h2),
            Table(
                [[route_drawing(day), route_block(day)]],
                colWidths=[96 * mm, 140 * mm],
                style=TableStyle([
                    ("VALIGN", (0, 0), (-1, -1), "TOP"),
                    ("LEFTPADDING", (0, 0), (-1, -1), 0),
                    ("RIGHTPADDING", (0, 0), (-1, -1), 6),
                    ("TOPPADDING", (0, 0), (-1, -1), 0),
                    ("BOTTOMPADDING", (0, 0), (-1, -1), 2),
                ]),
            ),
            Paragraph("Today's Theme", h2),
            p(today_theme(day)),
            Paragraph("時刻表", h2),
            pdf_table(
                [["Time", "Type", "Place", "Detail", "Duration"]] + [[i.time, i.type, i.place, short_text(i.detail, 76), i.duration] for i in day.timeline],
                [24 * mm, 22 * mm, 42 * mm, 132 * mm, 26 * mm],
                size=7.2,
            ),
        ]
        if rec_rows:
            story += [
                Paragraph("おすすめランチ・夕食", h2),
                pdf_table(
                    [["Type", "Name", "Rate", "Budget", "Popular"]] + rec_rows,
                    [34 * mm, 70 * mm, 22 * mm, 46 * mm, 58 * mm],
                    header=colors.HexColor("#F5EEDC"),
                    size=7.8,
                ),
            ]
        if tip_rows:
            story += [
                Paragraph("Today's Tips", h2),
                pdf_table([["Tips"]] + tip_rows, [238 * mm], header=colors.HexColor("#EEF6E8"), size=8.0),
            ]
        if spot_rows:
            story += [
                Paragraph("より道スポット・駐車場・Google Map", h2),
                pdf_table(
                    [["Spot", "Stay", "Parking", "GoogleMap"]] + spot_rows,
                    [54 * mm, 28 * mm, 56 * mm, 108 * mm],
                    header=colors.HexColor("#E2F1F8"),
                    size=6.8,
                ),
            ]
        story.append(PageBreak())

    doc.build(story)


def convert_pptx_to_pdf_if_possible() -> bool:
    soffice = shutil.which("libreoffice") or shutil.which("soffice")
    if not soffice:
        return False
    try:
        subprocess.run(
            [soffice, "--headless", "--convert-to", "pdf", "--outdir", str(OUTPUT_DIR), str(PPTX_PATH)],
            check=True,
            timeout=120,
        )
        converted = OUTPUT_DIR / "hokkaido-family-travel-guide.pdf"
        return converted.exists()
    except Exception:
        return False


def main():
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    PLACEHOLDER_DIR.mkdir(parents=True, exist_ok=True)
    days = load_days()
    if not days:
        raise SystemExit("No itinerary day files found under itinerary/days")
    build_pptx(days)
    if not convert_pptx_to_pdf_if_possible():
        build_pdf(days)
    print(f"Generated: {PPTX_PATH}")
    print(f"Generated: {PDF_PATH}")


if __name__ == "__main__":
    main()
